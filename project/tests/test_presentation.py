"""Tests for the reproducible stakeholder presentation build."""

from __future__ import annotations

import json
import shutil
from pathlib import Path

import pytest
from pptx import Presentation

from src.presentation import (
    SLIDES,
    build_presentation,
    date_span_years,
    format_iso_date,
    sample_rows_with_last,
)

PROJECT_ROOT = Path(__file__).resolve().parents[1]


def test_sample_rows_with_last_retains_newest_observation() -> None:
    rows = list(range(11))

    sampled = sample_rows_with_last(rows, 4)

    assert sampled == [0, 4, 8, 10]
    assert sampled[-1] == rows[-1]


@pytest.mark.parametrize("stride", [-3, 0, 1])
def test_sample_rows_with_last_stride_at_most_one_is_passthrough(stride: int) -> None:
    rows = [{"value": 1}, {"value": 2}]

    assert sample_rows_with_last(rows, stride) is rows


def test_sample_rows_with_last_single_row_is_passthrough() -> None:
    rows = [{"value": 1}]

    assert sample_rows_with_last(rows, 7) is rows


def test_sample_rows_with_last_does_not_duplicate_sampled_final_row() -> None:
    rows = list(range(7))

    assert sample_rows_with_last(rows, 3) == [0, 3, 6]


def test_format_iso_date() -> None:
    assert format_iso_date("2026-08-17") == "August 17, 2026"


def test_date_span_years_uses_calendar_anniversaries() -> None:
    assert date_span_years("2020-02-29", "2021-02-28") == pytest.approx(1.0)
    assert date_span_years("2020-01-01", "2020-07-02") == pytest.approx(0.5)
    assert date_span_years("2021-02-28", "2020-02-29") == pytest.approx(-1.0)


def test_build_presentation_writes_readable_authored_deck(tmp_path: Path) -> None:
    project_root = _copy_presentation_inputs(tmp_path)

    output_path = build_presentation(project_root)

    assert output_path == project_root / "reports/stakeholder_presentation.pptx"
    assert output_path.is_file()
    presentation = Presentation(output_path)
    assert presentation.core_properties.author == "Paritosh Dwivedi"
    assert presentation.core_properties.last_modified_by == "Paritosh Dwivedi"
    assert len(presentation.slides) == len(SLIDES)
    assert all(
        "[Sources]" in slide.notes_slide.notes_text_frame.text for slide in presentation.slides
    )


def _copy_presentation_inputs(tmp_path: Path) -> Path:
    destination = tmp_path / "project"
    metrics = json.loads((PROJECT_ROOT / "reports/metrics.json").read_text(encoding="utf-8"))
    raw_relative = Path(metrics["data"]["raw_file"])
    relative_paths = (
        Path("reports/metrics.json"),
        raw_relative,
        raw_relative.with_suffix(".manifest.json"),
        Path("reports/model_predictions.csv"),
        Path("reports/risk_threshold_sensitivity.csv"),
        Path("reports/feature_window_sensitivity.csv"),
        Path("reports/final_summary.md"),
        Path("README.md"),
        Path("docs/assumptions_and_risks.md"),
        Path("docs/decision_log.md"),
    )
    for relative_path in relative_paths:
        target = destination / relative_path
        target.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(PROJECT_ROOT / relative_path, target)
    return destination
