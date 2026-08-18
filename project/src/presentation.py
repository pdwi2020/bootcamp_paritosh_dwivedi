"""Build the editable stakeholder PowerPoint for the ETF risk monitor."""

from __future__ import annotations

import csv
import hashlib
import json
import math
import re
from collections.abc import Callable
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_DATA_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_MARKER_STYLE,
)
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

AUTHOR = "Paritosh Dwivedi"
OUTPUT_RELATIVE_PATH = Path("reports/stakeholder_presentation.pptx")

# Public declaration used by tests and reviewers to confirm the deck contract.
SLIDES = (
    "Title",
    "Decision question",
    "Validated history",
    "Leakage controls",
    "Regression evidence",
    "Classification evidence",
    "Robustness diagnostics",
    "Current signal",
    "Assumptions and risks",
    "Decision",
)

INK = "111111"
MUTED = "5F6368"
PANEL = "EDEDED"
RULE = "B8BCC4"
BLUE = "3D8DFF"
LIGHT_BLUE = "6DCBF4"
PALE_BLUE = "EAF5FB"
RED = "D92D20"
WHITE = "FFFFFF"
FONT_FAMILY = "Arial"


def sample_rows_with_last(rows: list, stride: int) -> list:
    """Sample every ``stride``-th row while always retaining the newest row."""

    if stride <= 1 or len(rows) <= 1:
        return rows
    sampled = rows[::stride]
    if sampled[-1] is not rows[-1]:
        sampled.append(rows[-1])
    return sampled


def format_iso_date(iso_date: str) -> str:
    """Format an ISO calendar date for stakeholder-facing prose."""

    parsed = date.fromisoformat(iso_date)
    return f"{parsed.strftime('%B')} {parsed.day}, {parsed.year}"


def date_span_years(start_iso: str, end_iso: str) -> float:
    """Return the signed calendar span, including the fractional anniversary year."""

    start = date.fromisoformat(start_iso)
    end = date.fromisoformat(end_iso)
    if start == end:
        return 0.0
    if end < start:
        return -date_span_years(end_iso, start_iso)

    whole_years = end.year - start.year
    anniversary = _anniversary(start, start.year + whole_years)
    if anniversary > end:
        whole_years -= 1
        anniversary = _anniversary(start, start.year + whole_years)
    next_anniversary = _anniversary(start, start.year + whole_years + 1)
    fraction = (end - anniversary).days / (next_anniversary - anniversary).days
    return whole_years + fraction


def _anniversary(value: date, year: int) -> date:
    """Move a date to ``year``, treating February 28 as the leap-day anniversary."""

    try:
        return value.replace(year=year)
    except ValueError:
        return value.replace(year=year, day=28)


def px(value: float) -> int:
    """Convert the reference deck's 96-DPI pixel coordinates to PowerPoint units."""

    return Inches(value / 96)


@dataclass(frozen=True)
class DeckInputs:
    """Validated source material required by the slide builders."""

    project_root: Path
    metrics: dict[str, Any]
    raw_relative_path: Path
    raw_rows: list[dict[str, Any]]
    manifest: dict[str, Any]
    predictions: list[dict[str, str]]
    threshold_sensitivity: list[dict[str, str]]
    feature_sensitivity: list[dict[str, str]]
    documents: dict[str, str]
    review_cadence: str
    risk_effects: dict[str, str]
    missing_required_values: int
    duplicate_dates: int


def build_presentation(project_root: Path) -> Path:
    """Build the stakeholder deck from committed project artifacts and return its path."""

    root = Path(project_root).resolve(strict=True)
    if not root.is_dir():
        raise NotADirectoryError(f"Project root is not a directory: {root}")
    inputs = _load_inputs(root)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    properties = presentation.core_properties
    properties.author = AUTHOR
    properties.last_modified_by = AUTHOR
    properties.title = "Weekly ETF Risk Monitor"
    properties.subject = "Stakeholder decision support for weekly SPY risk review"
    properties.keywords = "SPY, volatility, risk monitoring, FRE-GY 5040"

    builders: tuple[Callable[[Slide, DeckInputs], str], ...] = (
        _build_title_slide,
        _build_question_slide,
        _build_history_slide,
        _build_features_slide,
        _build_regression_slide,
        _build_classification_slide,
        _build_robustness_slide,
        _build_current_signal_slide,
        _build_risks_slide,
        _build_decision_slide,
    )
    if len(builders) != len(SLIDES):
        raise RuntimeError("Declared slide list and slide builders disagree")

    for number, (builder, note_entries) in enumerate(
        zip(builders, _speaker_note_sources(inputs), strict=True),
        start=1,
    ):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _set_background(slide, WHITE)
        footer_source = builder(slide, inputs)
        _add_footer(slide, number, footer_source)
        _add_speaker_notes(slide, note_entries)

    output_path = root / OUTPUT_RELATIVE_PATH
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return output_path


def _load_inputs(project_root: Path) -> DeckInputs:
    metrics_path = project_root / "reports/metrics.json"
    metrics = json.loads(metrics_path.read_text(encoding="utf-8"))
    if metrics.get("sole_author") != AUTHOR:
        raise ValueError(f"metrics.json must declare sole_author as {AUTHOR!r}")

    raw_relative_path, raw_path = _resolve_declared_raw_file(
        project_root,
        metrics.get("data", {}).get("raw_file"),
    )
    manifest_path = raw_path.with_suffix(".manifest.json")
    if not manifest_path.is_file():
        raise FileNotFoundError(f"Raw-data manifest is missing: {manifest_path}")
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    raw_rows, missing_required_values, duplicate_dates = _read_raw_rows(raw_path)
    _validate_raw_sources(metrics, raw_path, raw_rows, manifest)

    predictions = _read_csv(project_root / "reports/model_predictions.csv")
    threshold_sensitivity = _read_csv(project_root / "reports/risk_threshold_sensitivity.csv")
    feature_sensitivity = _read_csv(project_root / "reports/feature_window_sensitivity.csv")
    _validate_report_rows(metrics, predictions, threshold_sensitivity, feature_sensitivity)

    document_paths = {
        "final_summary": project_root / "reports/final_summary.md",
        "readme": project_root / "README.md",
        "assumptions": project_root / "docs/assumptions_and_risks.md",
        "decision_log": project_root / "docs/decision_log.md",
    }
    documents = {name: path.read_text(encoding="utf-8") for name, path in document_paths.items()}
    if any(not text.strip() for text in documents.values()):
        raise ValueError("Presentation source documents must not be empty")

    review_cadence = _extract_review_cadence(documents["readme"])
    risk_effects = _extract_risk_effects(documents["assumptions"])
    _validate_decision_documents(metrics, documents)

    return DeckInputs(
        project_root=project_root,
        metrics=metrics,
        raw_relative_path=raw_relative_path,
        raw_rows=raw_rows,
        manifest=manifest,
        predictions=predictions,
        threshold_sensitivity=threshold_sensitivity,
        feature_sensitivity=feature_sensitivity,
        documents=documents,
        review_cadence=review_cadence,
        risk_effects=risk_effects,
        missing_required_values=missing_required_values,
        duplicate_dates=duplicate_dates,
    )


def _resolve_declared_raw_file(project_root: Path, value: object) -> tuple[Path, Path]:
    if not isinstance(value, str) or not value.strip():
        raise ValueError("metrics.json data.raw_file must be a non-empty relative path")
    relative_path = Path(value)
    if relative_path.is_absolute():
        raise ValueError(f"Raw-data path must be relative to the project directory: {value}")

    resolved_root = project_root.resolve()
    resolved_path = (resolved_root / relative_path).resolve()
    try:
        resolved_path.relative_to(resolved_root)
    except ValueError as error:
        raise ValueError(f"Raw-data path escapes the project directory: {value}") from error
    if not resolved_path.is_file():
        raise FileNotFoundError(f"Raw-data file declared in metrics.json is missing: {value}")
    return relative_path, resolved_path


def _read_raw_rows(path: Path) -> tuple[list[dict[str, Any]], int, int]:
    with path.open(newline="", encoding="utf-8") as handle:
        reader = csv.DictReader(handle)
        required_columns = {"date", "adjusted_close"}
        if reader.fieldnames is None or not required_columns.issubset(reader.fieldnames):
            raise ValueError(f"Raw CSV lacks required columns {sorted(required_columns)}: {path}")
        records = list(reader)

    missing_required_values = sum(
        1
        for row in records
        for column in required_columns
        if row.get(column) is None or not str(row[column]).strip()
    )
    if missing_required_values:
        raise ValueError(f"Raw CSV contains {missing_required_values} missing required values")

    rows: list[dict[str, Any]] = []
    for row in records:
        iso_date = str(row["date"])
        date.fromisoformat(iso_date)
        try:
            price = float(row["adjusted_close"])
        except (TypeError, ValueError) as error:
            raise ValueError(f"Raw CSV has an invalid adjusted_close on {iso_date}") from error
        if not math.isfinite(price):
            raise ValueError(f"Raw CSV has a non-finite adjusted_close on {iso_date}")
        rows.append({"date": iso_date, "price": price})

    dates = [row["date"] for row in rows]
    duplicate_dates = len(dates) - len(set(dates))
    if duplicate_dates:
        raise ValueError(f"Raw CSV contains {duplicate_dates} duplicate dates")
    if dates != sorted(dates):
        raise ValueError("Raw CSV dates are not chronological")
    return rows, missing_required_values, duplicate_dates


def _validate_raw_sources(
    metrics: dict[str, Any],
    raw_path: Path,
    raw_rows: list[dict[str, Any]],
    manifest: dict[str, Any],
) -> None:
    if not raw_rows:
        raise ValueError("Raw CSV is empty")
    data = metrics["data"]
    expected = (data.get("rows_raw"), data.get("date_start"), data.get("date_end"))
    actual = (len(raw_rows), raw_rows[0]["date"], raw_rows[-1]["date"])
    if actual != expected:
        raise ValueError(
            "Raw chart data does not match rows_raw/date_start/date_end in metrics.json"
        )

    file_bytes = raw_path.read_bytes()
    checks = {
        "file": raw_path.name,
        "rows": len(raw_rows),
        "bytes": len(file_bytes),
        "sha256": hashlib.sha256(file_bytes).hexdigest(),
    }
    mismatches = [key for key, value in checks.items() if manifest.get(key) != value]
    if mismatches:
        raise ValueError(f"Raw-data manifest mismatch for: {', '.join(mismatches)}")


def _read_csv(path: Path) -> list[dict[str, str]]:
    if not path.is_file():
        raise FileNotFoundError(f"Presentation source is missing: {path}")
    with path.open(newline="", encoding="utf-8") as handle:
        reader = csv.DictReader(handle)
        if reader.fieldnames is None:
            raise ValueError(f"Presentation source has no CSV header: {path}")
        rows = list(reader)
    if not rows:
        raise ValueError(f"Presentation source has no data rows: {path}")
    return rows


def _validate_report_rows(
    metrics: dict[str, Any],
    predictions: list[dict[str, str]],
    threshold_sensitivity: list[dict[str, str]],
    feature_sensitivity: list[dict[str, str]],
) -> None:
    models = metrics["models"]
    if len(predictions) != models["test_rows"]:
        raise ValueError("Prediction row count does not match metrics.json")
    if predictions[0].get("date") != models["test_start"]:
        raise ValueError("Prediction start date does not match metrics.json")
    if predictions[-1].get("date") != models["test_end"]:
        raise ValueError("Prediction end date does not match metrics.json")

    threshold_metrics = metrics.get("sensitivity", {}).get("risk_threshold", [])
    feature_metrics = metrics.get("sensitivity", {}).get("feature_windows", [])
    if len(threshold_sensitivity) != len(threshold_metrics):
        raise ValueError("Risk-threshold sensitivity CSV does not match metrics.json")
    if len(feature_sensitivity) != len(feature_metrics):
        raise ValueError("Feature-window sensitivity CSV does not match metrics.json")


def _extract_review_cadence(readme: str) -> str:
    match = re.search(r"^\*\*Decision cadence:\*\*\s*(.+?)\s*$", readme, flags=re.MULTILINE)
    if match is None:
        raise ValueError("README.md does not declare the decision cadence")
    return match.group(1)


def _extract_risk_effects(document: str) -> dict[str, str]:
    rows: dict[str, str] = {}
    for line in document.splitlines():
        if not line.startswith("|"):
            continue
        cells = [cell.strip() for cell in line.strip().strip("|").split("|")]
        if len(cells) != 3 or cells[0] in {"Risk", "---"} or set(cells[0]) == {"-"}:
            continue
        rows[cells[0].casefold()] = cells[1].rstrip(".")

    required = {
        "regime change",
        "false reassurance",
        "score miscalibration",
        "tail underprediction",
    }
    missing = sorted(required.difference(rows))
    if missing:
        raise ValueError(f"Assumptions and risks document is missing required risks: {missing}")
    return rows


def _validate_decision_documents(metrics: dict[str, Any], documents: dict[str, str]) -> None:
    snapshot = metrics["latest_risk_snapshot"]
    decision_language = snapshot["decision_language"]
    if decision_language not in documents["final_summary"]:
        raise ValueError("final_summary.md does not match the latest decision language")

    rule = re.search(
        r"Use score\s*>=\s*([0-9]+(?:\.[0-9]+)?)% as the sole elevated-risk trigger",
        documents["decision_log"],
    )
    if rule is None:
        raise ValueError("decision_log.md does not declare the elevated-risk score rule")
    documented_cutoff = float(rule.group(1)) / 100
    if not math.isclose(documented_cutoff, float(snapshot["risk_score_cutoff"])):
        raise ValueError("decision_log.md score rule does not match the latest snapshot")


def _speaker_note_sources(inputs: DeckInputs) -> tuple[tuple[str, ...], ...]:
    raw_source = f"project/{inputs.raw_relative_path.as_posix()}"
    return (
        (raw_source, "project/reports/metrics.json"),
        ("project/README.md", "project/docs/decision_log.md"),
        (
            f"project/{inputs.raw_relative_path.with_suffix('.manifest.json').as_posix()}",
            "project/reports/metrics.json",
        ),
        (
            "project/src/features.py",
            "project/src/modeling.py",
            "project/docs/decision_log.md",
        ),
        ("project/reports/metrics.json", "project/reports/model_predictions.csv"),
        ("project/reports/metrics.json", "project/reports/model_predictions.csv"),
        (
            "project/reports/metrics.json",
            "project/reports/risk_threshold_sensitivity.csv",
            "project/reports/feature_window_sensitivity.csv",
        ),
        ("project/reports/metrics.json", "project/reports/final_summary.md"),
        (
            "project/docs/assumptions_and_risks.md",
            "project/docs/decision_log.md",
        ),
        (
            "project/reports/final_summary.md",
            "project/reports/metrics.json",
            "project/README.md",
        ),
    )


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def _set_background(slide: Slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_text(
    slide: Slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: float = 20,
    bold: bool = False,
    color: str = INK,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    name: str | None = None,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> Any:
    shape = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    if name:
        shape.name = name
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.vertical_anchor = vertical_anchor
    lines = text.split("\n")
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = alignment
        paragraph.font.name = FONT_FAMILY
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = _rgb(color)
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)
    return shape


def _add_rect(
    slide: Slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: str = PANEL,
    line: str | None = RULE,
) -> Any:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        px(left),
        px(top),
        px(width),
        px(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(0.75)
    return shape


def _add_title(slide: Slide, message: str) -> None:
    _add_text(
        slide,
        message,
        52,
        36,
        1176,
        66,
        font_size=35 if len(message) > 44 else 40,
        bold=True,
        name="slide-title",
    )
    _add_rect(slide, 52, 108, 1176, 1, fill=RULE, line=None)


def _add_footer(slide: Slide, number: int, source: str) -> None:
    _add_text(
        slide,
        source,
        52,
        674,
        1080,
        22,
        font_size=12,
        color=MUTED,
        name="source-footer",
    )
    _add_text(
        slide,
        str(number),
        1170,
        674,
        58,
        22,
        font_size=12,
        color=MUTED,
        alignment=PP_ALIGN.RIGHT,
        name="slide-number",
    )


def _add_speaker_notes(slide: Slide, entries: tuple[str, ...]) -> None:
    slide.notes_slide.notes_text_frame.text = (
        "[Sources]\n" + "\n".join(f"- {entry}" for entry in entries) + "\n[/Sources]"
    )


def _add_metric_callout(
    slide: Slide,
    x: float,
    y: float,
    label: str,
    value: str,
    explanation: str,
    accent: str = BLUE,
) -> None:
    _add_rect(slide, x, y, 260, 142, fill=PANEL, line=None)
    _add_rect(slide, x, y, 8, 142, fill=accent, line=None)
    _add_text(slide, value, x + 26, y + 18, 210, 50, font_size=34, bold=True)
    _add_text(slide, label, x + 26, y + 67, 210, 28, font_size=18, bold=True)
    _add_text(
        slide,
        explanation,
        x + 26,
        y + 99,
        210,
        34,
        font_size=15,
        color=MUTED,
    )


def _sample_for_chart(rows: list[dict[str, Any]], target_points: int) -> list[dict[str, Any]]:
    stride = max(1, math.ceil(len(rows) / target_points))
    return sample_rows_with_last(rows, stride)


def _add_line_chart(
    slide: Slide,
    rows: list[dict[str, Any]],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    background: str = WHITE,
    show_markers: bool = False,
) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = [row["date"][:7] for row in rows]
    chart_data.add_series("SPY adjusted close", [row["price"] for row in rows])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        px(left),
        px(top),
        px(width),
        px(height),
        chart_data,
    ).chart
    chart.has_title = False
    chart.has_legend = False
    del background
    series = chart.series[0]
    series.format.line.color.rgb = _rgb(BLUE)
    series.format.line.width = Pt(2.5)
    series.marker.style = XL_MARKER_STYLE.CIRCLE if show_markers else XL_MARKER_STYLE.NONE
    if show_markers:
        series.marker.size = 4
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = _rgb(BLUE)
        series.marker.format.line.color.rgb = _rgb(BLUE)
    _style_axes(chart, value_number_format="$0")


def _add_column_chart(
    slide: Slide,
    categories: list[str],
    series_rows: list[tuple[str, list[float], str]],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    maximum_scale: float | None = None,
    number_format: str = "0.0",
    legend: bool = True,
    show_values: bool = True,
    point_colors: dict[int, str] | None = None,
) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values, _ in series_rows:
        chart_data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        px(left),
        px(top),
        px(width),
        px(height),
        chart_data,
    ).chart
    chart.has_title = False
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = FONT_FAMILY
        chart.legend.font.size = Pt(12)
        chart.legend.font.color.rgb = _rgb(MUTED)
    _style_axes(chart, value_number_format=number_format, maximum_scale=maximum_scale)

    plot = chart.plots[0]
    plot.gap_width = 70
    if show_values:
        plot.has_data_labels = True
        labels = plot.data_labels
        labels.show_value = True
        labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        labels.number_format = number_format
        labels.font.name = FONT_FAMILY
        labels.font.size = Pt(11)
        labels.font.color.rgb = _rgb(INK)

    for series, (_, _, color) in zip(chart.series, series_rows, strict=True):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _rgb(color)
        series.format.line.fill.background()
    if point_colors:
        for index, color in point_colors.items():
            point = chart.series[0].points[index]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(color)
            point.format.line.fill.background()


def _style_axes(
    chart: Any,
    *,
    value_number_format: str,
    maximum_scale: float | None = None,
) -> None:
    category_axis = chart.category_axis
    category_axis.tick_labels.font.name = FONT_FAMILY
    category_axis.tick_labels.font.size = Pt(11)
    category_axis.tick_labels.font.color.rgb = _rgb(MUTED)
    category_axis.format.line.color.rgb = _rgb(RULE)

    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    if maximum_scale is not None:
        value_axis.maximum_scale = maximum_scale
    value_axis.tick_labels.number_format = value_number_format
    value_axis.tick_labels.font.name = FONT_FAMILY
    value_axis.tick_labels.font.size = Pt(10)
    value_axis.tick_labels.font.color.rgb = _rgb(MUTED)
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = _rgb(PANEL)
    value_axis.major_gridlines.format.line.width = Pt(0.75)
    value_axis.format.line.color.rgb = _rgb(RULE)


def _build_title_slide(slide: Slide, inputs: DeckInputs) -> str:
    metrics = inputs.metrics
    display_end_date = format_iso_date(metrics["data"]["date_end"])
    recent_session_count = min(756, len(inputs.raw_rows))
    recent_rows = _sample_for_chart(inputs.raw_rows[-recent_session_count:], 20)

    _add_text(
        slide,
        "FRE-GY 5040 · APPLIED FINANCIAL ENGINEERING",
        52,
        42,
        560,
        30,
        font_size=16,
        bold=True,
        color=BLUE,
    )
    _add_text(
        slide,
        "Weekly ETF\nRisk Monitor",
        52,
        154,
        520,
        178,
        font_size=64,
        bold=True,
        name="deck-title",
    )
    _add_text(
        slide,
        "A reproducible SPY signal for the portfolio manager's weekly risk review",
        52,
        362,
        510,
        92,
        font_size=25,
        color=MUTED,
    )
    _add_text(
        slide,
        f"{metrics['sole_author']} · Sole author\nData through {display_end_date}",
        52,
        544,
        480,
        62,
        font_size=18,
    )
    _add_rect(slide, 640, 42, 588, 588, fill=PALE_BLUE, line=None)
    _add_line_chart(
        slide,
        recent_rows,
        674,
        100,
        520,
        454,
        background=PALE_BLUE,
    )
    _add_text(
        slide,
        f"Recent {recent_session_count:,} validated trading sessions",
        700,
        566,
        468,
        30,
        font_size=15,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
    )
    return "Source: immutable SPY raw snapshot and project metrics"


def _build_question_slide(slide: Slide, inputs: DeckInputs) -> str:
    del inputs
    _add_title(slide, "One weekly question drives the entire pipeline")
    _add_text(
        slide,
        "Should the portfolio manager maintain SPY exposure—or investigate reducing it "
        "because near-term risk appears elevated?",
        58,
        164,
        548,
        270,
        font_size=34,
        bold=True,
    )
    _add_rect(slide, 654, 154, 574, 440, fill=PANEL, line=None)
    _add_text(
        slide,
        "THE USEFUL ANSWER",
        692,
        186,
        480,
        28,
        font_size=18,
        bold=True,
        color=BLUE,
    )
    rows = (
        ("1", "Forecast", "Next-five-session annualized volatility", 246),
        ("2", "Warning", "Class-weighted risk score—not a literal probability", 350),
        ("3", "Interpretation", "Action language plus assumptions and limits", 454),
    )
    for number, label, detail, top in rows:
        _add_text(slide, number, 694, top, 44, 44, font_size=30, bold=True, color=BLUE)
        _add_text(slide, label, 754, top, 220, 32, font_size=24, bold=True)
        _add_text(slide, detail, 754, top + 36, 400, 54, font_size=18, color=MUTED)
    return "Decision support—not automated trading advice"


def _build_history_slide(slide: Slide, inputs: DeckInputs) -> str:
    data = inputs.metrics["data"]
    observation_count = len(inputs.raw_rows)
    history_years = date_span_years(data["date_start"], data["date_end"])
    all_rows = _sample_for_chart(inputs.raw_rows, 22)
    quality_issues = inputs.missing_required_values + inputs.duplicate_dates

    _add_title(slide, f"{observation_count:,} validated trading days anchor the analysis")
    _add_text(
        slide,
        "SPY adjusted close · full validated history",
        52,
        120,
        760,
        24,
        font_size=16,
        color=MUTED,
        name="chart-window-label",
    )
    _add_line_chart(slide, all_rows, 52, 150, 760, 480, show_markers=True)
    _add_metric_callout(
        slide,
        860,
        150,
        "OBSERVATIONS",
        f"{observation_count:,}",
        "Daily SPY rows",
        BLUE,
    )
    _add_metric_callout(
        slide,
        860,
        314,
        "HISTORY",
        f"{history_years:.1f} yrs",
        f"{data['date_start']} to {data['date_end']}",
        LIGHT_BLUE,
    )
    _add_metric_callout(
        slide,
        860,
        478,
        "QUALITY",
        str(quality_issues),
        "Missing required values or duplicates",
        BLUE,
    )
    provider = inputs.manifest.get("provider", "declared provider")
    return f"Source: {provider} fallback; immutable CSV + SHA-256 manifest"


def _build_features_slide(slide: Slide, inputs: DeckInputs) -> str:
    models = inputs.metrics["models"]
    _add_title(slide, "Every feature is available before the forecast window")
    xs = (52, 348, 644, 940)
    labels = (
        ("1 · SNAPSHOT", "Immutable daily prices and volume"),
        ("2 · TRAILING SIGNALS", "Returns, volatility, drawdown, volume"),
        ("3 · FUTURE TARGET", "Next five trading days of realized volatility"),
        ("4 · LATER TEST", "Chronological out-of-sample evaluation"),
    )
    for index, ((heading, detail), x) in enumerate(zip(labels, xs, strict=True)):
        _add_rect(slide, x, 210, 242, 292, fill=PALE_BLUE if index == 2 else PANEL, line=None)
        _add_text(
            slide,
            heading,
            x + 24,
            242,
            194,
            56,
            font_size=20,
            bold=True,
            color=BLUE if index == 2 else INK,
        )
        _add_text(slide, detail, x + 24, 324, 194, 106, font_size=20, color=MUTED)
        if index < len(labels) - 1:
            _add_text(
                slide,
                "→",
                x + 252,
                318,
                34,
                44,
                font_size=30,
                bold=True,
                color=BLUE,
                alignment=PP_ALIGN.CENTER,
            )
    _add_text(
        slide,
        "Leakage control: threshold and outlier parameters use training data only; "
        f"{models['embargo_rows']} sessions are embargoed.",
        52,
        548,
        1120,
        62,
        font_size=24,
        bold=True,
    )
    return (
        "Method: five-session horizon, trailing features, "
        f"purged chronological {round(models['train_rows'] / (models['train_rows'] + models['test_rows']) * 100):.0f}/"
        f"{round(models['test_rows'] / (models['train_rows'] + models['test_rows']) * 100):.0f} split"
    )


def _build_regression_slide(slide: Slide, inputs: DeckInputs) -> str:
    metrics = inputs.metrics
    regression = metrics["models"]["regression"]
    improvement = regression["ridge_mae_improvement_vs_recent"] * 100
    _add_title(
        slide,
        f"Ridge cut forecast MAE by {improvement:.1f}% versus recent volatility",
    )
    _add_text(
        slide,
        "Annualized volatility error (percentage points)",
        52,
        126,
        520,
        28,
        font_size=15,
        color=MUTED,
    )
    _add_column_chart(
        slide,
        ["MAE", "RMSE"],
        [
            (
                "Ridge",
                [regression["ridge"][metric] * 100 for metric in ("mae", "rmse")],
                BLUE,
            ),
            (
                "Recent-vol baseline",
                [
                    regression["recent_volatility_baseline"][metric] * 100
                    for metric in ("mae", "rmse")
                ],
                LIGHT_BLUE,
            ),
            (
                "Historical mean",
                [
                    regression["historical_mean_baseline"][metric] * 100
                    for metric in ("mae", "rmse")
                ],
                RULE,
            ),
        ],
        52,
        160,
        720,
        458,
    )
    _add_metric_callout(
        slide,
        842,
        160,
        "RIDGE MAE",
        f"{regression['ridge']['mae'] * 100:.2f} pp",
        "Annualized-vol error",
        BLUE,
    )
    _add_metric_callout(
        slide,
        842,
        324,
        "R-SQUARED",
        f"{regression['ridge']['r2']:.3f}",
        "Out-of-sample",
        LIGHT_BLUE,
    )
    _add_metric_callout(
        slide,
        842,
        488,
        "IMPROVEMENT",
        f"{improvement:.1f}%",
        "Versus recent-vol MAE",
        BLUE,
    )
    models = metrics["models"]
    return f"Test period: {models['test_start']} to {models['test_end']}"


def _build_classification_slide(slide: Slide, inputs: DeckInputs) -> str:
    summary = inputs.metrics["models"]["classification"]
    classification = summary["logistic"]
    recall = classification["recall"] * 100
    _add_title(slide, f"The warning model catches {recall:.0f}% of elevated-risk windows")
    categories = ["Balanced accuracy", "Precision", "Recall", "ROC AUC"]
    values = [
        classification["balanced_accuracy"] * 100,
        classification["precision"] * 100,
        classification["recall"] * 100,
        classification["roc_auc"] * 100,
    ]
    _add_column_chart(
        slide,
        categories,
        [("Logistic risk model", values, BLUE)],
        52,
        164,
        728,
        450,
        maximum_scale=100,
        number_format='0.0"%"',
        legend=False,
    )
    _add_rect(slide, 834, 166, 394, 188, fill=PALE_BLUE, line=None)
    _add_text(
        slide,
        str(summary["elevated_windows_caught"]),
        868,
        194,
        100,
        58,
        font_size=42,
        bold=True,
        color=BLUE,
    )
    _add_text(
        slide,
        "elevated-risk windows caught",
        868,
        257,
        310,
        36,
        font_size=22,
        bold=True,
    )
    _add_text(
        slide,
        f"{summary['elevated_windows_missed']} elevated-risk windows were missed",
        868,
        306,
        310,
        28,
        font_size=17,
        color=MUTED,
    )
    _add_text(
        slide,
        "Why balanced accuracy?",
        834,
        405,
        350,
        34,
        font_size=25,
        bold=True,
    )
    _add_text(
        slide,
        f"Elevated-risk windows are only {summary['test_elevated_rate'] * 100:.1f}% of "
        "the holdout. The prior baseline predicts only the majority class and catches zero "
        "elevated-risk windows.",
        834,
        452,
        382,
        126,
        font_size=18,
        color=MUTED,
    )
    return "Decision emphasis: elevated-risk recall and balanced accuracy"


def _build_robustness_slide(slide: Slide, inputs: DeckInputs) -> str:
    diagnostics = inputs.metrics["models"]["diagnostics"]
    non_overlap = diagnostics["non_overlapping_windows"]
    non_overlap_recall = [row["recall"] * 100 for row in non_overlap]
    calendar_year = diagnostics["calendar_year"]
    calendar_low = min(calendar_year, key=lambda row: row["recall"])
    calendar_high = max(calendar_year, key=lambda row: row["recall"])
    walk_forward = diagnostics["walk_forward"]["aggregate"]
    residuals = diagnostics["residuals"]

    _add_title(slide, "Robustness checks temper confidence—not the conclusion")
    _add_text(
        slide,
        "Overlapping targets, changing regimes, and tail errors all need explicit checks.",
        52,
        126,
        1120,
        42,
        font_size=22,
        color=MUTED,
    )
    rows = (
        (
            f"{min(non_overlap_recall):.0f}–{max(non_overlap_recall):.0f}%",
            "NON-OVERLAPPING RECALL",
            f"{len(non_overlap)} offsets sample every fifth forecast window; no offset reverses "
            "the warning model's usefulness.",
        ),
        (
            f"{calendar_low['recall'] * 100:.0f}–{calendar_high['recall'] * 100:.0f}%",
            "CALENDAR-YEAR RECALL",
            f"Recall varies from {calendar_low['year']} to {calendar_high['year']}, confirming "
            "that market regimes materially affect detection.",
        ),
        (
            f"{walk_forward['recall'] * 100:.1f}%",
            "WALK-FORWARD RECALL",
            f"{len(diagnostics['walk_forward']['folds'])} expanding folds with a "
            f"{inputs.metrics['models']['embargo_rows']}-session embargo; balanced accuracy "
            f"{walk_forward['balanced_accuracy'] * 100:.1f}%.",
        ),
        (
            f"+{residuals['top_actual_vol_decile_mean_actual_minus_predicted'] * 100:.1f} pp",
            "TOP-DECILE UNDERPREDICTION",
            "Ridge understates the most volatile windows, so tail-risk decisions still require "
            "human judgment.",
        ),
    )
    for index, (value, label, detail) in enumerate(rows):
        top = 176 + index * 113
        _add_rect(
            slide,
            52,
            top,
            1176,
            96,
            fill=PALE_BLUE if index == 2 else PANEL,
            line=None,
        )
        _add_text(
            slide,
            value,
            78,
            top + 20,
            224,
            52,
            font_size=31,
            bold=True,
            color=BLUE if index == 2 else INK,
        )
        _add_text(slide, label, 334, top + 15, 390, 30, font_size=19, bold=True)
        _add_text(slide, detail, 334, top + 49, 842, 42, font_size=16, color=MUTED)

    quantiles = [float(row["risk_quantile"]) * 100 for row in inputs.threshold_sensitivity]
    return (
        f"Also tested: {min(quantiles):.0f}%–{max(quantiles):.0f}% risk quantiles, "
        f"{len(inputs.feature_sensitivity)} feature-window specifications, and outlier-flag ablation"
    )


def _build_current_signal_slide(slide: Slide, inputs: DeckInputs) -> str:
    snapshot = inputs.metrics["latest_risk_snapshot"]
    classification = str(snapshot["risk_classification"])
    _add_title(slide, f"Current signal is {classification.lower()}—not risk-free")
    signal_color = RED if classification.casefold() == "elevated" else BLUE
    _add_text(
        slide,
        classification.upper(),
        52,
        144,
        520,
        104,
        font_size=76,
        bold=True,
        color=signal_color,
    )
    _add_text(
        slide,
        f"As of {snapshot['as_of_date']}",
        58,
        254,
        360,
        34,
        font_size=20,
        color=MUTED,
    )
    _add_text(
        slide,
        snapshot["decision_language"],
        58,
        344,
        500,
        130,
        font_size=28,
        bold=True,
    )
    _add_text(
        slide,
        f"A {classification.lower()} signal means the model does not detect elevated relative "
        "risk under its historical definition.",
        58,
        508,
        500,
        82,
        font_size=19,
        color=MUTED,
    )
    _add_text(
        slide,
        "Annualized volatility (%)",
        642,
        122,
        300,
        28,
        font_size=15,
        color=MUTED,
    )
    volatility_values = [
        snapshot["rolling_vol_20"] * 100,
        snapshot["predicted_next_five_day_vol"] * 100,
        snapshot["risk_threshold_annualized_vol"] * 100,
    ]
    chart_max = math.ceil(max(volatility_values) * 1.25)
    _add_column_chart(
        slide,
        ["20-day actual", "Five-session forecast", "Risk threshold"],
        [("Annualized volatility", volatility_values, BLUE)],
        642,
        150,
        586,
        430,
        maximum_scale=chart_max,
        number_format="0.0",
        legend=False,
        point_colors={2: RED},
    )
    _add_text(
        slide,
        f"Risk score: {snapshot['elevated_risk_score'] * 100:.1f}% · elevated at "
        f"{snapshot['risk_score_cutoff'] * 100:.0f}%",
        672,
        592,
        530,
        34,
        font_size=22,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    return "The signal supports review; it does not authorize a trade"


def _build_risks_slide(slide: Slide, inputs: DeckInputs) -> str:
    _add_title(slide, "Four risks define how the signal should be used")
    risk_copy = (
        (
            "regime change",
            "Historical relationships can fail when market structure shifts.",
        ),
        ("false reassurance", "Normal means not elevated by this model—not safe."),
        (
            "score miscalibration",
            "Class weighting makes the output a risk score—not a literal probability.",
        ),
        (
            "tail underprediction",
            "Ridge understates the most volatile forecast windows.",
        ),
    )
    positions = ((52, 172), (654, 172), (52, 414), (654, 414))
    for index, ((risk, reference_copy), (left, top)) in enumerate(
        zip(risk_copy, positions, strict=True)
    ):
        # The risk register supplies the required items; the visible wording follows the
        # stakeholder reference deck while the source effect is validated at load time.
        if not inputs.risk_effects[risk]:
            raise ValueError(f"Risk register effect is empty: {risk}")
        _add_rect(
            slide,
            left,
            top,
            574,
            180,
            fill=PALE_BLUE if index == 1 else PANEL,
            line=None,
        )
        _add_text(
            slide,
            risk.upper(),
            left + 28,
            top + 28,
            510,
            34,
            font_size=22,
            bold=True,
            color=BLUE if index == 1 else INK,
        )
        _add_text(
            slide,
            reference_copy,
            left + 28,
            top + 82,
            510,
            66,
            font_size=20,
            color=MUTED,
        )
    return (
        "Mitigations: hash validation, purged tests, walk-forward checks, tail diagnostics, "
        "explicit caveats"
    )


def _build_decision_slide(slide: Slide, inputs: DeckInputs) -> str:
    snapshot = inputs.metrics["latest_risk_snapshot"]
    _add_text(slide, "DECISION", 52, 56, 220, 28, font_size=17, bold=True, color=BLUE)
    _add_text(
        slide,
        snapshot["decision_language"],
        52,
        126,
        920,
        196,
        font_size=46,
        bold=True,
    )
    _add_rect(slide, 52, 356, 1176, 1, fill=RULE, line=None)
    cutoff = snapshot["risk_score_cutoff"] * 100
    _add_text(
        slide,
        "Investigate exposure or hedging when the class-weighted elevated-risk score reaches "
        f"{cutoff:.0f}%. Use the volatility forecast as supporting context.",
        52,
        408,
        830,
        116,
        font_size=27,
        color=MUTED,
    )
    _add_rect(slide, 966, 408, 262, 142, fill=PALE_BLUE, line=None)
    _add_text(
        slide,
        "NEXT REVIEW",
        994,
        434,
        210,
        28,
        font_size=17,
        bold=True,
        color=BLUE,
        alignment=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        inputs.review_cadence,
        994,
        476,
        210,
        48,
        font_size=34,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        f"{inputs.metrics['sole_author']} · Sole author",
        52,
        626,
        440,
        30,
        font_size=18,
    )
    return "Decision support—not a guarantee or automated trade"
